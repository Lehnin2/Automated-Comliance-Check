"""
Enhanced PowerPoint Extractor with Intelligent Parser
Extracts raw data, then uses LLM to extract ONLY relevant raw text (NO JUDGMENTS)
Uses LLM with automatic fallback (TokenFactory -> Gemini)
"""

import json
import re
from datetime import datetime
from pathlib import Path
from typing import Dict, List

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Import path utilities for centralized path management
from path_utils import DOCUMENTS_DIR, UPLOADS_DIR, RESULTS_DIR

# Import LLM Manager with fallback support
from llm_manager import llm_manager


class RawExtractor:
    """Pure extraction from PowerPoint only"""
    
    def __init__(self, pptx_path: str):
        self.pptx_path = Path(pptx_path)
        self.pres = Presentation(pptx_path)
        
    def extract(self) -> Dict:
        """Extract everything from PowerPoint"""
        print(f"  📄 Extracting {len(self.pres.slides)} slides...")
        
        pages = []
        for idx, slide in enumerate(self.pres.slides, 1):
            pages.append(self._extract_slide(slide, idx))
            print(f"    ✓ Page {idx}")
        
        return {
            "document": {
                "filename": self.pptx_path.name,
                "total_pages": len(pages),
                "extracted_at": datetime.now().isoformat()[:19]
            },
            "pages": pages
        }
    
    def _extract_slide(self, slide, page_num: int) -> Dict:
        """Extract one slide - raw data only"""
        
        texts = []
        
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame") or not shape.text.strip():
                continue
            
            # Get position
            try:
                y_pos = round(shape.top.inches / self.pres.slide_height.inches, 3) if hasattr(shape, 'top') else None
                x_pos = round(shape.left.inches, 2) if hasattr(shape, 'left') else None
            except:
                y_pos = None
                x_pos = None
            
            # Extract text with formatting
            runs = []
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    
                    try:
                        font_size = run.font.size.pt if run.font.size else None
                        font_name = run.font.name
                        is_bold = run.font.bold if run.font.bold is not None else None
                        is_italic = run.font.italic if run.font.italic is not None else None
                        
                        color = None
                        if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                            rgb = run.font.color.rgb
                            color = f"#{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"
                    except:
                        font_size = None
                        font_name = None
                        is_bold = None
                        is_italic = None
                        color = None
                    
                    runs.append({
                        "text": run.text,
                        "font_size": font_size,
                        "font_name": font_name,
                        "bold": is_bold,
                        "italic": is_italic,
                        "color": color
                    })
            
            if runs:
                texts.append({
                    "full_text": shape.text.strip(),
                    "position_y": y_pos,
                    "position_x": x_pos,
                    "formatting": runs
                })
        
        # Extract tables
        tables = []
        for shape in slide.shapes:
            if shape.has_table:
                table_data = []
                for row in shape.table.rows:
                    row_cells = []
                    for cell in row.cells:
                        is_bold = False
                        try:
                            if cell.text_frame and cell.text_frame.paragraphs:
                                para = cell.text_frame.paragraphs[0]
                                if para.runs and para.runs[0].font.bold:
                                    is_bold = True
                        except:
                            pass
                        
                        row_cells.append({
                            "text": cell.text.strip(),
                            "bold": is_bold
                        })
                    table_data.append(row_cells)
                
                tables.append({
                    "rows": len(shape.table.rows),
                    "cols": len(shape.table.columns),
                    "data": table_data
                })
        
        # Count visuals
        chart_count = sum(1 for s in slide.shapes if s.has_chart)
        image_count = sum(1 for s in slide.shapes if s.shape_type == MSO_SHAPE_TYPE.PICTURE)
        
        return {
            "page": page_num,
            "texts": texts,
            "tables": tables,
            "charts": chart_count,
            "images": image_count
        }


class IntelligentParser:
    """Uses LLM to extract ONLY relevant raw text - NO JUDGMENTS"""
    
    def __init__(self, api_key: str = None):
        # Use llm_manager with automatic fallback
        self.llm = llm_manager
    
    def parse_for_compliance(self, raw_extraction: Dict) -> Dict:
        """Parse extracted data and extract only relevant raw text"""
        
        print(f"\n  🤖 Parsing with LLM for compliance relevance...")
        
        compliance_structure = {
            "document_info": {},
            "cover_page": {},
            "slide_2_disclaimers": {},
            "content_pages": [],
            "end_page": {},
            "performance_sections": [],
            "esg_content": [],
            "fund_characteristics": {},
            "risk_warnings": [],
            "all_bold_text": [],
            "all_sources_citations": []
        }
        
        pages = raw_extraction["pages"]
        
        # Parse Cover Page (Page 1)
        if pages:
            print(f"    • Parsing Cover Page...")
            compliance_structure["cover_page"] = self._parse_cover_page(pages[0])
        
        # Parse Slide 2 (Disclaimers)
        if len(pages) > 1:
            print(f"    • Parsing Slide 2 (Disclaimers)...")
            compliance_structure["slide_2_disclaimers"] = self._parse_slide_2(pages[1])
        
        # Parse remaining content pages
        print(f"    • Parsing Content Pages...")
        for page in pages[2:]:
            parsed_page = self._parse_content_page(page)
            if parsed_page:
                compliance_structure["content_pages"].append(parsed_page)
        
        # Parse last page
        if pages:
            print(f"    • Parsing End Page...")
            compliance_structure["end_page"] = self._parse_end_page(pages[-1])
        
        # Extract cross-page elements (raw text only)
        print(f"    • Extracting cross-page elements...")
        compliance_structure["performance_sections"] = self._extract_performance_sections(pages)
        compliance_structure["esg_content"] = self._extract_esg_content(pages)
        compliance_structure["all_bold_text"] = self._extract_all_bold_text(pages)
        compliance_structure["all_sources_citations"] = self._extract_sources(pages)
        
        # Add document metadata
        compliance_structure["document_info"] = {
            "filename": raw_extraction["document"]["filename"],
            "total_pages": raw_extraction["document"]["total_pages"],
            "extraction_timestamp": raw_extraction["document"]["extracted_at"]
        }
        
        return compliance_structure
    
    def _call_llm(self, prompt: str, context: str) -> str:
        """Call LLM for parsing with automatic fallback"""
        result = self.llm.call_llm(
            system_prompt=prompt,
            user_prompt=context,
            temperature=0.1,
            max_tokens=2000
        )
        if not result:
            print(f"      ⚠️  LLM call failed: {self.llm.last_error}")
            return ""
        return result
    
    def _parse_cover_page(self, page: Dict) -> Dict:
        """Extract ONLY relevant raw text from cover page - NO JUDGMENTS"""
        
        all_text = "\n".join([t["full_text"] for t in page["texts"]])
        
        prompt = """You are extracting raw text for compliance checking. Extract ONLY the exact text (word-for-word) for these elements:
1. fund_name: Extract the complete fund name exactly as written
2. date_mention: Extract any date/month/year text exactly as shown
3. promotional_mention: Extract exact text if "document promotionnel" or similar appears
4. audience_mention: Extract exact text mentioning "retail", "professionnel", "professional", "non-professional"
5. confidentiality_mention: Extract exact text if "do not disclose", "confidential" or similar appears
6. pre_commercial_warning: Extract the FULL red/bold warning text if present (about fund not yet approved)
7. client_specific_mention: Extract exact text if a specific client name is mentioned

Return ONLY as JSON with these keys. Use null if element not found. Return EXACT TEXT ONLY, no interpretation.
Example: {"fund_name": "ODDO BHF US Equity Active UCITS ETF", "date_mention": "October 2025", ...}"""

        llm_response = self._call_llm(prompt, all_text)
        
        try:
            parsed = json.loads(llm_response)
        except:
            parsed = {}
        
        # Add raw text sections for verification
        top_texts = [t["full_text"] for t in page["texts"] if t["position_y"] and t["position_y"] < 0.3]
        
        return {
            "extracted_elements": parsed,
            "full_top_section_text": top_texts,
            "visual_count": {"images": page["images"], "charts": page["charts"]}
        }
    
    def _parse_slide_2(self, page: Dict) -> Dict:
        """Extract raw text from slide 2 - disclaimers and mandatory mentions"""
        
        all_text = "\n".join([t["full_text"] for t in page["texts"]])
        
        prompt = """Extract ONLY the exact raw text (word-for-word) for:
1. disclaimer_text: Extract the full disclaimer paragraph (retail or professional)
2. risk_profile_text: Extract complete risk profile description
3. countries_text: Extract text mentioning commercialization countries
4. sri_mention: Extract exact text showing SRI value and explanation
5. bold_risk_warnings: Extract all bold text about risks (exact text)

Return as JSON. Use null if not found. NO interpretation, only exact text extraction."""

        llm_response = self._call_llm(prompt, all_text)
        
        try:
            parsed = json.loads(llm_response)
        except:
            parsed = {}
        
        # Extract all bold text manually
        bold_texts = []
        for text_block in page["texts"]:
            for run in text_block["formatting"]:
                if run["bold"] and len(run["text"]) > 20:
                    bold_texts.append({
                        "text": run["text"],
                        "font_size": run["font_size"]
                    })
        
        return {
            "extracted_elements": parsed,
            "all_bold_text_found": bold_texts
        }
    
    def _parse_content_page(self, page: Dict) -> Dict:
        """Extract relevant raw text from content pages"""
        
        all_text = "\n".join([t["full_text"] for t in page["texts"]])
        
        if len(all_text.strip()) < 50:
            return None
        
        prompt = """Extract ONLY exact raw text for:
1. performance_text: If page shows performance data, extract ALL text mentioning percentages, returns, periods (1Y, 3Y, etc.)
2. morningstar_text: Extract exact text mentioning Morningstar rating and date
3. stock_value_mentions: Extract exact sentences mentioning specific company/stock names
4. esg_text: Extract exact text mentioning ESG, sustainability, SFDR
5. team_text: Extract text presenting management team members
6. source_citations: Extract all source mentions (bottom of page or footnotes)

Return as JSON. Use null if not found. EXACT TEXT ONLY."""

        llm_response = self._call_llm(prompt, all_text[:3000])
        
        try:
            parsed = json.loads(llm_response)
        except:
            parsed = {}
        
        # Extract table data if present
        table_texts = []
        if page["tables"]:
            for table in page["tables"]:
                table_texts.append({
                    "dimensions": f"{table['rows']}x{table['cols']}",
                    "first_row": [cell["text"] for cell in table["data"][0]] if table["data"] else []
                })
        
        return {
            "page_number": page["page"],
            "extracted_text": parsed,
            "tables": table_texts,
            "visual_count": {"charts": page["charts"], "images": page["images"]}
        }
    
    def _parse_end_page(self, page: Dict) -> Dict:
        """Extract legal mentions from last page - raw text only"""
        
        all_text = "\n".join([t["full_text"] for t in page["texts"]])
        
        prompt = """Extract ONLY exact raw text for:
1. sgp_legal_mention: Extract the complete legal text about the management company (société de gestion)
2. characteristics_table_text: If fund characteristics table present, extract key-value pairs
3. glossary_text: Extract glossary content if present
4. validation_responsible_text: Extract text mentioning who validated the document

Return as JSON. Use null if not found. EXACT TEXT ONLY."""

        llm_response = self._call_llm(prompt, all_text)
        
        try:
            parsed = json.loads(llm_response)
        except:
            parsed = {}
        
        # Bottom section text
        bottom_texts = [t["full_text"] for t in page["texts"] if t["position_y"] and t["position_y"] > 0.7]
        
        return {
            "extracted_elements": parsed,
            "full_bottom_section_text": bottom_texts
        }
    
    def _extract_performance_sections(self, pages: List[Dict]) -> List[Dict]:
        """Extract ALL text from performance-related sections"""
        
        performance_sections = []
        
        for page in pages:
            all_text = " ".join([t["full_text"] for t in page["texts"]])
            
            # Check for performance keywords
            perf_keywords = ["performance", "rendement", "ytd", "year to date", "annualisé", "cumulé", "%"]
            if any(kw in all_text.lower() for kw in perf_keywords):
                
                prompt = """Extract ONLY exact raw text related to performance:
1. performance_values_text: Extract all text showing performance numbers with time periods
2. benchmark_comparison_text: Extract text comparing to benchmark
3. performance_disclaimer_text: Extract complete disclaimer about past performance
4. chart_title_text: Extract chart title if performance chart present

Return as JSON. EXACT TEXT ONLY."""

                llm_response = self._call_llm(prompt, all_text[:2000])
                
                try:
                    parsed = json.loads(llm_response)
                    if any(parsed.values()):  # Only add if something was found
                        parsed["page_number"] = page["page"]
                        performance_sections.append(parsed)
                except:
                    pass
        
        return performance_sections
    
    def _extract_esg_content(self, pages: List[Dict]) -> List[Dict]:
        """Extract ALL ESG-related text"""
        
        esg_content = []
        
        for page in pages:
            all_text = " ".join([t["full_text"] for t in page["texts"]])
            
            esg_keywords = ["esg", "environnement", "social", "gouvernance", "durable", "responsable", "sfdr", "sustainability"]
            if any(kw in all_text.lower() for kw in esg_keywords):
                
                prompt = """Extract ONLY exact raw text mentioning ESG/sustainability:
1. esg_approach_text: Extract text describing ESG approach or methodology
2. sfdr_classification_text: Extract text about SFDR Article 6/8/9
3. exclusion_criteria_text: Extract text about ESG exclusions
4. esg_integration_text: Extract text about how ESG is integrated

Return as JSON. EXACT TEXT ONLY."""

                llm_response = self._call_llm(prompt, all_text[:2000])
                
                try:
                    parsed = json.loads(llm_response)
                    if any(parsed.values()):
                        parsed["page_number"] = page["page"]
                        esg_content.append(parsed)
                except:
                    pass
        
        return esg_content
    
    def _extract_all_bold_text(self, pages: List[Dict]) -> List[Dict]:
        """Extract ALL bold text across document (for risk warnings check)"""
        
        all_bold = []
        
        for page in pages:
            for text_block in page["texts"]:
                for run in text_block["formatting"]:
                    if run["bold"] and len(run["text"].strip()) > 15:
                        all_bold.append({
                            "page": page["page"],
                            "text": run["text"],
                            "font_size": run["font_size"],
                            "color": run["color"]
                        })
        
        return all_bold
    
    def _extract_sources(self, pages: List[Dict]) -> List[Dict]:
        """Extract ALL source citations across document"""
        
        all_sources = []
        
        for page in pages:
            all_text = " ".join([t["full_text"] for t in page["texts"]])
            
            prompt = """Extract ALL source citations from this page.
Look for: "Source:", "Sources:", footnote references, data attributions.
Return as JSON array: [{"source_text": "exact citation text"}]
EXACT TEXT ONLY."""

            llm_response = self._call_llm(prompt, all_text[:2000])
            
            try:
                parsed = json.loads(llm_response)
                if parsed:
                    all_sources.append({
                        "page": page["page"],
                        "citations": parsed
                    })
            except:
                pass
        
        return all_sources


def process_document(
    pptx_path: str,
    metadata_file: str = None,
    output_path: str = None,
    use_llm_parser: bool = True,
    api_key: str = "YOUR_API_KEY"
) -> Dict:
    """Main processing with intelligent parsing (raw text extraction only)"""
    
    print("\n" + "="*80)
    print(f"🚀 SMART EXTRACTION: {Path(pptx_path).name}")
    print("="*80)
    
    # Load user metadata
    user_metadata = None
    if metadata_file and Path(metadata_file).exists():
        print(f"\n  📥 Loading user metadata...")
        with open(metadata_file, 'r', encoding='utf-8') as f:
            user_metadata = json.load(f)
        print(f"    ✓ Loaded")
    
    # Step 1: Raw extraction
    print(f"\n  📤 Step 1: Raw extraction from PowerPoint...")
    extractor = RawExtractor(pptx_path)
    raw_extracted = extractor.extract()
    
    raw_json = json.dumps(raw_extracted, indent=2)
    raw_lines = len(raw_json.split('\n'))
    raw_size_kb = len(raw_json.encode('utf-8')) / 1024
    
    print(f"    ✓ Raw extraction: {raw_lines} lines, {raw_size_kb:.1f} KB")
    
    # Step 2: Intelligent parsing
    if use_llm_parser:
        print(f"\n  🧠 Step 2: Extracting compliance-relevant raw text...")
        parser = IntelligentParser(api_key=api_key)
        compliance_data = parser.parse_for_compliance(raw_extracted)
        
        # Add user metadata
        compliance_data["user_metadata"] = user_metadata
        
        parsed_json = json.dumps(compliance_data, indent=2, ensure_ascii=False)
        parsed_lines = len(parsed_json.split('\n'))
        parsed_size_kb = len(parsed_json.encode('utf-8')) / 1024
        
        reduction_pct = ((raw_lines - parsed_lines) / raw_lines) * 100
        
        print(f"    ✓ Extraction complete")
        print(f"    ✓ Reduced: {raw_lines} → {parsed_lines} lines ({reduction_pct:.1f}% reduction)")
        print(f"    ✓ Size: {raw_size_kb:.1f} KB → {parsed_size_kb:.1f} KB")
        
        final_output = compliance_data
    else:
        final_output = raw_extracted
        parsed_lines = raw_lines
    
    # Save
    if output_path:
        output_file = Path(output_path)
        output_file.parent.mkdir(parents=True, exist_ok=True)
        
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(final_output, f, indent=2, ensure_ascii=False)
        
        print(f"\n  💾 Saved: {output_path}")
    
    # Summary
    print("\n" + "="*80)
    print("📊 EXTRACTION SUMMARY")
    print("="*80)
    print(f"  • Input: {Path(pptx_path).name}")
    print(f"  • Total pages: {raw_extracted['document']['total_pages']}")
    
    if use_llm_parser:
        print(f"\n  • Raw extraction: {raw_lines} lines")
        print(f"  • Compliance-ready: {parsed_lines} lines")
        print(f"  • Reduction: {reduction_pct:.1f}%")
        
        print(f"\n  📋 Extracted elements (RAW TEXT ONLY):")
        print(f"     - Cover page elements")
        print(f"     - Slide 2 disclaimers")
        print(f"     - Content pages: {len(compliance_data['content_pages'])}")
        print(f"     - Performance sections: {len(compliance_data['performance_sections'])}")
        print(f"     - ESG content: {len(compliance_data['esg_content'])}")
        print(f"     - Bold text items: {len(compliance_data['all_bold_text'])}")
        print(f"     - Source citations: {len(compliance_data['all_sources_citations'])}")
    
    print("="*80)
    print("\n✅ READY FOR COMPLIANCE CHECKING")
    print("📝 Output contains ONLY raw text - NO JUDGMENTS")
    print("Your friend's code will make all compliance decisions!")
    
    return final_output


if __name__ == "__main__":
    
    result = process_document(
        pptx_path="XXX-PRS-GB-ODDO BHF US Equity Active ETF-20250630_6PN.pptx",
        metadata_file="metadata2.json",
        output_path="output/compliance_ready_extraction2.json",
        use_llm_parser=True,
        api_key="sk-5298061f39b44388b6539c2d504e6a18"
    )
    
    print(f"\n🎯 Your friend gets:")
    print(f"   ✓ Only relevant raw text extracts")
    print(f"   ✓ No true/false judgments")
    print(f"   ✓ Organized by compliance categories")
    print(f"   ✓ Ready for his rule engine!")