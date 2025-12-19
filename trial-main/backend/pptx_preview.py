"""
PowerPoint Preview Generator
Converts PPTX slides to images for web preview
Uses PowerPoint COM on Windows for high-fidelity rendering
"""

from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import io
import base64
from pathlib import Path
import json
import os
import tempfile
import shutil
from logger_config import logger


def convert_pptx_to_images_com(pptx_path: str, output_dir: str) -> list:
    """
    Convert PPTX to images using PowerPoint COM (Windows only)
    Returns list of image file paths
    """
    try:
        import comtypes.client
        
        # Initialize COM
        powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
        powerpoint.Visible = 1  # Required for export
        
        # Open presentation
        abs_path = os.path.abspath(pptx_path)
        presentation = powerpoint.Presentations.Open(abs_path, WithWindow=False)
        
        image_paths = []
        
        try:
            # Export each slide as PNG
            for i, slide in enumerate(presentation.Slides, start=1):
                output_path = os.path.join(output_dir, f"slide_{i}.png")
                slide.Export(output_path, "PNG", 1280, 720)
                image_paths.append(output_path)
                logger.debug(f"Exported slide {i} to {output_path}")
        finally:
            presentation.Close()
            powerpoint.Quit()
        
        return image_paths
        
    except Exception as e:
        logger.warning(f"COM conversion failed: {e}")
        return []


def convert_pptx_to_images_libreoffice(pptx_path: str, output_dir: str) -> list:
    """
    Convert PPTX to images using LibreOffice (cross-platform)
    Returns list of image file paths
    """
    try:
        import subprocess
        
        # First convert to PDF
        pdf_dir = tempfile.mkdtemp()
        
        result = subprocess.run([
            'soffice',
            '--headless',
            '--convert-to', 'pdf',
            '--outdir', pdf_dir,
            pptx_path
        ], capture_output=True, timeout=60)
        
        if result.returncode != 0:
            logger.warning(f"LibreOffice conversion failed: {result.stderr}")
            return []
        
        # Find the PDF
        pdf_name = Path(pptx_path).stem + '.pdf'
        pdf_path = os.path.join(pdf_dir, pdf_name)
        
        if not os.path.exists(pdf_path):
            return []
        
        # Convert PDF pages to images using pdf2image
        try:
            from pdf2image import convert_from_path
            images = convert_from_path(pdf_path, dpi=150)
            
            image_paths = []
            for i, img in enumerate(images, start=1):
                output_path = os.path.join(output_dir, f"slide_{i}.png")
                img.save(output_path, 'PNG')
                image_paths.append(output_path)
            
            return image_paths
        except ImportError:
            logger.warning("pdf2image not installed")
            return []
        finally:
            shutil.rmtree(pdf_dir, ignore_errors=True)
            
    except Exception as e:
        logger.warning(f"LibreOffice conversion failed: {e}")
        return []


def render_slide_to_image_pil(slide, width=960, height=540) -> str:
    """
    Fallback: Render a slide to an image using PIL
    Creates a styled visual representation of the slide
    """
    try:
        # Create image with gradient-like background
        img = Image.new('RGB', (width, height), color='#f8f9fa')
        draw = ImageDraw.Draw(img)
        
        # Draw a subtle border
        draw.rectangle([0, 0, width-1, height-1], outline='#dee2e6', width=2)
        
        # Try to load fonts
        try:
            title_font = ImageFont.truetype("arial.ttf", 28)
            body_font = ImageFont.truetype("arial.ttf", 16)
            small_font = ImageFont.truetype("arial.ttf", 12)
        except:
            try:
                title_font = ImageFont.truetype("C:/Windows/Fonts/arial.ttf", 28)
                body_font = ImageFont.truetype("C:/Windows/Fonts/arial.ttf", 16)
                small_font = ImageFont.truetype("C:/Windows/Fonts/arial.ttf", 12)
            except:
                title_font = ImageFont.load_default()
                body_font = ImageFont.load_default()
                small_font = ImageFont.load_default()
        
        y_position = 40
        title_found = False
        
        # Draw shapes/text
        for shape in slide.shapes:
            if hasattr(shape, 'text') and shape.text.strip():
                text = shape.text.strip()
                
                # Determine if it's a title
                is_title = False
                if hasattr(shape, 'name') and 'Title' in shape.name:
                    is_title = True
                elif hasattr(shape, 'is_placeholder') and shape.is_placeholder:
                    try:
                        if hasattr(shape, 'placeholder_format') and shape.placeholder_format.type == 1:
                            is_title = True
                    except:
                        pass
                
                # First large text is likely title
                if not title_found and len(text) < 100:
                    is_title = True
                    title_found = True
                
                if is_title:
                    font = title_font
                    color = '#1e3a5f'  # Dark blue for titles
                    # Draw title with underline
                    lines = wrap_text(text, font, width - 80, draw)
                    for line in lines[:2]:  # Max 2 lines for title
                        if y_position < height - 50:
                            draw.text((40, y_position), line, fill=color, font=font)
                            y_position += 35
                    # Draw underline
                    draw.line([(40, y_position), (width - 40, y_position)], fill='#3498db', width=2)
                    y_position += 25
                else:
                    font = body_font
                    color = '#2c3e50'  # Dark gray for body
                    
                    # Wrap and draw text
                    lines = wrap_text(text, font, width - 100, draw)
                    for line in lines[:8]:  # Limit lines
                        if y_position < height - 40:
                            # Add bullet point for body text
                            draw.text((50, y_position), "•", fill='#3498db', font=body_font)
                            draw.text((70, y_position), line, fill=color, font=font)
                            y_position += 24
                    
                    y_position += 10  # Space between elements
        
        # Convert to base64
        buffer = io.BytesIO()
        img.save(buffer, format='PNG', quality=95)
        img_str = base64.b64encode(buffer.getvalue()).decode()
        
        return f"data:image/png;base64,{img_str}"
        
    except Exception as e:
        logger.error(f"Error rendering slide to image: {e}")
        return None


def wrap_text(text, font, max_width, draw):
    """Wrap text to fit within max_width"""
    words = text.split()
    lines = []
    current_line = []
    
    for word in words:
        test_line = ' '.join(current_line + [word])
        try:
            bbox = draw.textbbox((0, 0), test_line, font=font)
            width = bbox[2] - bbox[0]
        except:
            width = len(test_line) * 8  # Fallback estimation
        
        if width <= max_width:
            current_line.append(word)
        else:
            if current_line:
                lines.append(' '.join(current_line))
            current_line = [word]
    
    if current_line:
        lines.append(' '.join(current_line))
    
    return lines


def render_slide_to_image(slide, width=960, height=540) -> str:
    """
    Main function to render slide to image
    Tries COM first (best quality), then falls back to PIL
    """
    # For individual slides, use PIL rendering
    # COM/LibreOffice is used for batch conversion in extract_slide_thumbnails
    return render_slide_to_image_pil(slide, width, height)


def pptx_to_images(pptx_path: str, output_dir: str = None) -> list:
    """
    Convert PowerPoint slides to base64 encoded images
    
    Args:
        pptx_path: Path to PPTX file
        output_dir: Optional directory to save images
        
    Returns:
        List of dicts with slide info and base64 image data
    """
    try:
        from pptx.util import Inches
        import comtypes.client
        import os
        
        # Convert PPTX to images using COM (Windows only)
        # For cross-platform, we'll extract slide content as JSON
        prs = Presentation(pptx_path)
        
        slides_data = []
        
        for idx, slide in enumerate(prs.slides, start=1):
            slide_info = {
                'slide_number': idx,
                'shapes': [],
                'notes': ''
            }
            
            # Extract text and shapes
            for shape in slide.shapes:
                shape_data = {
                    'type': shape.shape_type.name if hasattr(shape.shape_type, 'name') else 'UNKNOWN',
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height
                }
                
                # Extract text
                if hasattr(shape, 'text'):
                    shape_data['text'] = shape.text
                
                # Extract image if present
                if hasattr(shape, 'image'):
                    try:
                        image_bytes = shape.image.blob
                        image_b64 = base64.b64encode(image_bytes).decode('utf-8')
                        shape_data['image'] = f"data:image/png;base64,{image_b64}"
                    except:
                        pass
                
                slide_info['shapes'].append(shape_data)
            
            # Extract notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    slide_info['notes'] = notes_slide.notes_text_frame.text
            
            slides_data.append(slide_info)
        
        return slides_data
        
    except Exception as e:
        print(f"Error converting PPTX: {e}")
        return []


def extract_slide_thumbnails(pptx_path: str) -> list:
    """
    Extract slide content as structured data for preview with images
    
    Tries to use PowerPoint COM (Windows) for high-fidelity images,
    falls back to LibreOffice, then to PIL rendering.
    
    Returns:
        List of slide data with text content, layout info, and rendered images
    """
    logger.info(f"Extracting slides from: {pptx_path}")
    
    try:
        prs = Presentation(pptx_path)
        slides = []
        num_slides = len(prs.slides)
        
        logger.info(f"Found {num_slides} slides")
        
        # Try to get high-fidelity images using COM or LibreOffice
        high_fidelity_images = {}
        temp_dir = None
        
        try:
            temp_dir = tempfile.mkdtemp(prefix="pptx_preview_")
            
            # Try PowerPoint COM first (Windows with Office installed)
            image_paths = convert_pptx_to_images_com(pptx_path, temp_dir)
            
            if not image_paths:
                # Try LibreOffice as fallback
                logger.info("Trying LibreOffice conversion...")
                image_paths = convert_pptx_to_images_libreoffice(pptx_path, temp_dir)
            
            if image_paths:
                logger.info(f"High-fidelity conversion successful: {len(image_paths)} images")
                for i, img_path in enumerate(image_paths, start=1):
                    try:
                        with open(img_path, 'rb') as f:
                            img_data = f.read()
                            img_b64 = base64.b64encode(img_data).decode()
                            high_fidelity_images[i] = f"data:image/png;base64,{img_b64}"
                    except Exception as e:
                        logger.warning(f"Failed to read image {i}: {e}")
            else:
                logger.info("No high-fidelity conversion available, using PIL fallback")
                
        except Exception as e:
            logger.warning(f"High-fidelity conversion failed: {e}")
        
        # Process each slide
        for idx, slide in enumerate(prs.slides, start=1):
            logger.debug(f"Processing slide {idx}")
            
            slide_data = {
                'slide_number': idx,
                'title': '',
                'content': [],
                'layout': slide.slide_layout.name if hasattr(slide.slide_layout, 'name') else 'Unknown',
                'image': None,  # Will contain base64 image
                'has_high_fidelity': False
            }
            
            # Use high-fidelity image if available, otherwise render with PIL
            if idx in high_fidelity_images:
                slide_data['image'] = high_fidelity_images[idx]
                slide_data['has_high_fidelity'] = True
                logger.debug(f"Slide {idx}: using high-fidelity image")
            else:
                try:
                    slide_data['image'] = render_slide_to_image(slide)
                    logger.debug(f"Slide {idx}: rendered with PIL")
                except Exception as e:
                    logger.error(f"Failed to render slide {idx}: {e}")
            
            # Extract all text content (always needed for compliance checking)
            for shape in slide.shapes:
                if hasattr(shape, 'text') and shape.text.strip():
                    text_info = {
                        'text': shape.text.strip(),
                        'is_title': shape.name.startswith('Title') if hasattr(shape, 'name') else False,
                        'font_size': None,
                        'position': {
                            'left': float(shape.left) if hasattr(shape, 'left') else 0,
                            'top': float(shape.top) if hasattr(shape, 'top') else 0,
                            'width': float(shape.width) if hasattr(shape, 'width') else 0,
                            'height': float(shape.height) if hasattr(shape, 'height') else 0
                        }
                    }
                    
                    # Try to get font size
                    if hasattr(shape, 'text_frame'):
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                if run.font.size:
                                    text_info['font_size'] = run.font.size.pt
                                    break
                            if text_info['font_size']:
                                break
                    
                    # Set title if it looks like one
                    if text_info['is_title'] or (text_info['font_size'] and text_info['font_size'] > 24):
                        slide_data['title'] = text_info['text']
                    
                    slide_data['content'].append(text_info)
            
            slides.append(slide_data)
            logger.debug(f"Slide {idx} extracted: {len(slide_data['content'])} elements")
        
        # Cleanup temp directory
        if temp_dir:
            try:
                shutil.rmtree(temp_dir, ignore_errors=True)
            except:
                pass
        
        logger.info(f"Successfully extracted {len(slides)} slides")
        return slides
        
    except Exception as e:
        logger.error(f"Error extracting slides: {e}")
        raise
