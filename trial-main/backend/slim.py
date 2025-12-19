import json
import re
import httpx
import asyncio
import time
from datetime import datetime
from typing import List, Dict, Any, Optional
from concurrent.futures import ThreadPoolExecutor, as_completed
from openai import OpenAI
from pydantic import BaseModel, Field
from pptx import Presentation
import os

# ==================== SETUP CLIENT ====================

http_client = httpx.Client(verify=False)

client = OpenAI(
    api_key="sk-97bf4a384b5d4b1cb2d3199ac6659917",  # Replace with your key
    base_url="https://tokenfactory.esprit.tn/api",
    http_client=http_client
)

# ==================== PYDANTIC MODELS ====================

class Position(BaseModel):
    page: int = Field(default=1)
    order: int = Field(default=1)
    paragraph: Optional[int] = None

class Claim(BaseModel):
    claim_id: str
    claim_text: str
    claim_type: str
    section_reference: str
    position: Position
    sources_cited: List[str] = []
    evidence_dates: List[str] = []
    is_qualified: bool = False

class Disclaimer(BaseModel):
    disclaimer_id: str
    disclaimer_text: str
    disclaimer_type: str
    section_reference: str
    position: Position
    font_size: Optional[int] = None
    is_bold: bool = False
    visibility: bool = True

class Fee(BaseModel):
    percentage: Optional[float] = None
    note: Optional[str] = None

class AUM(BaseModel):
    value: Optional[float] = None
    currency: str = "EUR"
    date: Optional[str] = None

class FundCharacteristics(BaseModel):
    inception_date: Optional[str] = None
    aum: AUM = Field(default_factory=AUM)
    benchmark: Optional[str] = None
    investment_horizon: Optional[str] = None
    risk_level: Optional[int] = None

class Fees(BaseModel):
    management_fee: Fee = Field(default_factory=Fee)
    subscription_fee: Fee = Field(default_factory=Fee)
    performance_fee: Fee = Field(default_factory=Fee)
    ter: Fee = Field(default_factory=Fee)

class KeyDataPoints(BaseModel):
    fund_characteristics: FundCharacteristics = Field(default_factory=FundCharacteristics)
    fees: Fees = Field(default_factory=Fees)

class DocumentMetadata(BaseModel):
    document_id: str
    document_name: str
    document_type: str
    creation_date: str
    last_updated: str
    validated_by: Optional[str] = None
    language: str = "EN"
    country: Optional[str] = None
    page_count: int = 1
    fund_name: str
    fund_isin: Optional[str] = None
    fund_type: str = "UCITS"

# ==================== PPTX EXTRACTION ====================

def extract_text_from_pptx(pptx_path: str) -> str:
    """Extract all text from PPTX file using LangExtract"""
    print(f"📖 Lecture du fichier PPTX: {pptx_path}")

    try:
        # Try using langextract first
        from langextract import extract
        print("🔧 Utilisant LangExtract...")

        text = extract(pptx_path)
        print(f"✅ {len(text)} caractères extraits avec LangExtract")
        return text
    except Exception as e:
        print(f"⚠️ LangExtract échoué: {e}")
        print("📖 Fallback: Utilisant python-pptx...")

        # Fallback to python-pptx
        try:
            prs = Presentation(pptx_path)
            full_text = []

            for slide_idx, slide in enumerate(prs.slides, 1):
                slide_text = f"\n--- SLIDE {slide_idx} ---\n"

                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        if shape.text.strip():
                            slide_text += shape.text + "\n"

                full_text.append(slide_text)

            text = "\n".join(full_text)
            print(f"✅ {len(prs.slides)} slides extraites avec python-pptx")
            return text
        except Exception as e2:
            print(f"❌ Erreur lecture PPTX: {e2}")
            return ""

# ==================== LLM HELPER FUNCTIONS ====================

def call_llm(prompt: str, max_tokens: int = 1000) -> str:
    """Call the LLM via OpenAI client"""
    try:
        response = client.chat.completions.create(
            model="hosted_vllm/Llama-3.1-70B-Instruct",
            messages=[
                {"role": "system", "content": "Tu es un expert en analyse de documents financiers. Réponds toujours en JSON valide."},
                {"role": "user", "content": prompt}
            ],
            temperature=0.3,
            max_tokens=max_tokens,
            top_p=0.9,
            frequency_penalty=0.0,
            presence_penalty=0.0
        )
        return response.choices[0].message.content
    except Exception as e:
        print(f"❌ Erreur LLM: {e}")
        return "{}"

def safe_json_parse(text: str) -> Dict:
    """Safely parse JSON from LLM response"""
    try:
        return json.loads(text)
    except:
        try:
            json_match = re.search(r'\{.*\}|\[.*\]', text, re.DOTALL)
            if json_match:
                return json.loads(json_match.group())
        except:
            pass
    return {}

# ==================== EXTRACTION CLASS WITH PARALLEL EXECUTION ====================

class DocumentParser:
    def __init__(self, num_workers: int = 4):
        self.document_id = f"doc_{datetime.now().strftime('%Y%m%d_%H%M%S')}"
        self.num_workers = num_workers
        self.extraction_times = {}

    def extract_metadata(self, document_text: str) -> DocumentMetadata:
        """Extract metadata from document"""
        prompt = f"""Analyse ce document financier et extrais les métadonnées clés en JSON.

Document (premiers 2000 chars):
{document_text[:2000]}

Extrais et retourne UNIQUEMENT un JSON avec:
- document_name: nom du fonds
- document_type: type de document
- creation_date: date de création (format YYYY-MM-DD, sinon aujourd'hui)
- language: langue (FR, EN, DE)
- page_count: nombre de pages (entier)
- fund_name: nom du fonds
- fund_isin: code ISIN (ou null)
- fund_type: type de fonds (UCITS, ETF, etc)

Réponds UNIQUEMENT avec le JSON."""

        response_text = call_llm(prompt, max_tokens=500)
        data = safe_json_parse(response_text)

        try:
            metadata = DocumentMetadata(
                document_id=self.document_id,
                document_name=data.get("document_name", "Unknown"),
                document_type=data.get("document_type", "fund_factsheet"),
                creation_date=data.get("creation_date", datetime.now().isoformat()[:10]),
                last_updated=datetime.now().isoformat()[:10],
                language=data.get("language", "EN"),
                page_count=int(data.get("page_count", 1)),
                fund_name=data.get("fund_name", ""),
                fund_isin=data.get("fund_isin"),
                fund_type=data.get("fund_type", "UCITS")
            )
            return metadata
        except Exception as e:
            print(f"⚠️ Erreur parsing metadata: {e}")
            return None

    def extract_claims(self, document_text: str) -> List[Claim]:
        """Extract investment claims from document"""
        prompt = f"""Identifie TOUTES les CLAIMS (affirmations/promesses) dans ce document financier.

Document (premiers 3000 chars):
{document_text[:3000]}

Pour CHAQUE claim trouvée, extrais:
- claim_text: le texte exact
- claim_type: "performance", "risk", "benefit", "feature", "strategy"
- sources_cited: liste des sources
- evidence_dates: années/dates mentionnées
- is_qualified: true si la claim a un disclaimer

Retourne un JSON array UNIQUEMENT. Minimum 3 claims."""

        response_text = call_llm(prompt, max_tokens=1500)
        data = safe_json_parse(response_text)

        claims = []
        if isinstance(data, list):
            for idx, claim_data in enumerate(data):
                try:
                    claim = Claim(
                        claim_id=f"claim_{idx+1}",
                        claim_text=claim_data.get("claim_text", ""),
                        claim_type=claim_data.get("claim_type", "feature"),
                        section_reference="general",
                        position=Position(page=1, order=idx+1),
                        sources_cited=claim_data.get("sources_cited", []),
                        evidence_dates=claim_data.get("evidence_dates", []),
                        is_qualified=claim_data.get("is_qualified", False)
                    )
                    claims.append(claim)
                except Exception as e:
                    print(f"⚠️ Erreur parsing claim {idx}: {e}")
        return claims

    def extract_disclaimers(self, document_text: str) -> List[Disclaimer]:
        """Extract disclaimers and warnings"""
        prompt = f"""Identifie TOUS les DISCLAIMERS et WARNINGS dans ce document.

Document (premiers 3000 chars):
{document_text[:3000]}

Pour CHAQUE disclaimer trouvé, extrais:
- disclaimer_text: le texte complet
- disclaimer_type: "risk_warning", "performance_caveat", "liability_limitation", "general_warning"
- is_bold: true si probablement en gras

Retourne un JSON array UNIQUEMENT. Minimum 2 disclaimers."""

        response_text = call_llm(prompt, max_tokens=1500)
        data = safe_json_parse(response_text)

        disclaimers = []
        if isinstance(data, list):
            for idx, disc_data in enumerate(data):
                try:
                    disclaimer = Disclaimer(
                        disclaimer_id=f"disclaimer_{idx+1}",
                        disclaimer_text=disc_data.get("disclaimer_text", ""),
                        disclaimer_type=disc_data.get("disclaimer_type", "general_warning"),
                        section_reference="disclaimers",
                        position=Position(page=1, order=idx+1),
                        is_bold=disc_data.get("is_bold", False),
                        visibility=True
                    )
                    disclaimers.append(disclaimer)
                except Exception as e:
                    print(f"⚠️ Erreur parsing disclaimer {idx}: {e}")
        return disclaimers

    def extract_key_data(self, document_text: str) -> KeyDataPoints:
        """Extract key financial data"""
        prompt = f"""Extrais les données financières clés de ce document.

Document (premiers 3000 chars):
{document_text[:3000]}

Extrais (ou mets null si non trouvé):
- inception_date: date de création (YYYY-MM-DD)
- aum_value: actifs sous gestion (nombre)
- aum_currency: devise (EUR, USD, GBP, etc)
- benchmark: indice de référence
- investment_horizon: horizon d'investissement
- risk_level: niveau de risque (1-7, entier)
- management_fee: frais de gestion (nombre en %)
- subscription_fee: frais de souscription (nombre en %)
- performance_fee: frais de performance (nombre en %)
- ter: Total Expense Ratio (nombre en %)

Retourne un JSON UNIQUEMENT."""

        response_text = call_llm(prompt, max_tokens=800)
        data = safe_json_parse(response_text)

        try:
            key_data = KeyDataPoints(
                fund_characteristics=FundCharacteristics(
                    inception_date=data.get("inception_date"),
                    aum=AUM(
                        value=data.get("aum_value"),
                        currency=data.get("aum_currency", "EUR"),
                        date=datetime.now().isoformat()[:10]
                    ),
                    benchmark=data.get("benchmark"),
                    investment_horizon=data.get("investment_horizon"),
                    risk_level=data.get("risk_level")
                ),
                fees=Fees(
                    management_fee=Fee(percentage=data.get("management_fee")),
                    subscription_fee=Fee(percentage=data.get("subscription_fee")),
                    performance_fee=Fee(percentage=data.get("performance_fee")),
                    ter=Fee(percentage=data.get("ter"))
                )
            )
            return key_data
        except Exception as e:
            print(f"⚠️ Erreur key_data: {e}")
            return KeyDataPoints()

    def parse_document_parallel(self, document_text: str) -> Dict[str, Any]:
        """Parse document with parallel execution"""
        print("\n⚡ Démarrage de l'extraction PARALLÈLE...\n")

        start_time = time.time()

        extraction_tasks = {
            "metadata": lambda: self.extract_metadata(document_text),
            "claims": lambda: self.extract_claims(document_text),
            "disclaimers": lambda: self.extract_disclaimers(document_text),
            "key_data": lambda: self.extract_key_data(document_text)
        }

        results = {}

        # Parallel execution with ThreadPoolExecutor
        with ThreadPoolExecutor(max_workers=self.num_workers) as executor:
            future_to_task = {
                executor.submit(task_func): task_name
                for task_name, task_func in extraction_tasks.items()
            }

            for future in as_completed(future_to_task):
                task_name = future_to_task[future]
                task_start = time.time()
                try:
                    results[task_name] = future.result()
                    task_time = time.time() - task_start
                    self.extraction_times[task_name] = task_time
                    print(f"✅ {task_name} terminé en {task_time:.2f}s")
                except Exception as e:
                    print(f"❌ Erreur dans {task_name}: {e}")
                    results[task_name] = None
                    self.extraction_times[task_name] = 0

        total_time = time.time() - start_time

        metadata = results.get("metadata")
        claims = results.get("claims", [])
        disclaimers = results.get("disclaimers", [])
        key_data = results.get("key_data")

        final_result = {
            "document_metadata": metadata.model_dump() if metadata else {},
            "claims": [claim.model_dump() for claim in claims],
            "disclaimers_and_warnings": [disc.model_dump() for disc in disclaimers],
            "key_data_points": key_data.model_dump() if key_data else {},
            "compliance_metadata": {
                "extraction_method": "automated_llm_parallel",
                "extraction_date": datetime.now().isoformat(),
                "confidence_score": 0.75,
                "flagged_for_review": False,
                "parallel_workers": self.num_workers
            },
            "performance_metrics": {
                "total_extraction_time": f"{total_time:.2f}s",
                "metadata_time": f"{self.extraction_times.get('metadata', 0):.2f}s",
                "claims_time": f"{self.extraction_times.get('claims', 0):.2f}s",
                "disclaimers_time": f"{self.extraction_times.get('disclaimers', 0):.2f}s",
                "key_data_time": f"{self.extraction_times.get('key_data', 0):.2f}s"
            }
        }
        return final_result

# ==================== MAIN ====================

def main():
    # Upload PPTX file
    from google.colab import files

    print("📤 Sélectionne ton fichier PPTX...\n")
    uploaded = files.upload()

    if not uploaded:
        print("❌ Aucun fichier uploadé")
        return

    pptx_filename = list(uploaded.keys())[0]
    print(f"✅ Fichier reçu: {pptx_filename}\n")

    # Extract text from PPTX
    document_text = extract_text_from_pptx(pptx_filename)

    if not document_text:
        print("❌ Erreur lors de l'extraction du texte")
        return

    print(f"📝 Texte extrait: {len(document_text)} caractères\n")

    # Parse with parallel execution
    parser = DocumentParser(num_workers=4)
    result = parser.parse_document_parallel(document_text)

    print("\n✅ Extraction terminée!\n")

    # Export to JSON
    output_filename = f"parsed_document_{datetime.now().strftime('%Y%m%d_%H%M%S')}.json"
    with open(output_filename, 'w', encoding='utf-8') as f:
        json.dump(result, f, indent=2, ensure_ascii=False)

    print(f"💾 Document exporté: {output_filename}")
    print(f"\n📊 Résumé:")
    print(f"   - Claims identifiées: {len(result['claims'])}")
    print(f"   - Disclaimers identifiés: {len(result['disclaimers_and_warnings'])}")
    print(f"   - Workers parallèles: {result['compliance_metadata']['parallel_workers']}")

    # Display performance metrics
    print(f"\n⏱️  PERFORMANCE METRICS:")
    perf = result['performance_metrics']
    print(f"   - Total time: {perf['total_extraction_time']}")
    print(f"   - Metadata: {perf['metadata_time']}")
    print(f"   - Claims: {perf['claims_time']}")
    print(f"   - Disclaimers: {perf['disclaimers_time']}")
    print(f"   - Key Data: {perf['key_data_time']}")

    # Display JSON
    print(f"\n{json.dumps(result, indent=2, ensure_ascii=False)}")

    # Download JSON
    files.download(output_filename)

# RUN
if __name__ == "__main__":
    main()