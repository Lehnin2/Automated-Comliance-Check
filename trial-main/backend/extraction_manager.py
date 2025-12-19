"""
Extraction Manager - Handles multiple extraction methods
Supports: 
  - MO: Current Method (python-pptx) - Basic extraction
  - FD: Fida Method (Gemini Multi-Agent) - AI-powered with LangGraph
  - SF: Safa Method (Groq) - Complete exhaustive extraction with caching
  - SL: Slim Method (TokenFactory) - Parallel extraction with claims/disclaimers
"""

import os
import json
from pathlib import Path
from datetime import datetime
from typing import Dict, Any, Optional
from pptx import Presentation
from dotenv import load_dotenv

# Import path utilities and load env
from path_utils import ENV_FILE
load_dotenv(str(ENV_FILE))

# Import logger
from logger_config import logger


class ExtractionManager:
    """Manages different extraction methods"""
    
    METHODS = {
        'MO': 'Standard (python-pptx)',
        'FD': 'Fida (Gemini Multi-Agent)',
        'SF': 'Safa (Groq Exhaustive)',
        'SL': 'Slim (TokenFactory Parallel)'
    }
    
    def __init__(self):
        # Reload env to ensure we have the latest values
        load_dotenv(str(ENV_FILE))
        self.gemini_api_key = os.environ.get('GEMINI_API_KEY', '')
        self.groq_api_key = os.environ.get('GROQ_API_KEY', '')
        self.tokenfactory_api_key = os.environ.get('TOKENFACTORY_API_KEY', '')
    
    def extract(self, pptx_path: str, method: str = 'MO', output_path: str = None, parallel_workers: int = 4, on_event: Optional[callable] = None) -> Dict[str, Any]:
        """
        Extract content from PPTX using specified method
        
        Args:
            pptx_path: Path to PowerPoint file
            method: 'MO', 'FD', 'SF', or 'SL'
            output_path: Optional path to save extracted JSON
            parallel_workers: Number of parallel workers for SL method (1-8)
            on_event: Optional callback to receive streaming events (FD only)
            
        Returns:
            Extracted data dictionary
        """
        logger.info(f"Starting extraction with method: {method}")
        
        if method == 'FD':
            return self._extract_fida(pptx_path, output_path, on_event)
        elif method == 'SF':
            return self._extract_safa(pptx_path, output_path)
        elif method == 'SL':
            return self._extract_slim(pptx_path, output_path, parallel_workers)
        else:
            return self._extract_mo(pptx_path, output_path)
    
    def _extract_mo(self, pptx_path: str, output_path: str = None) -> Dict[str, Any]:
        """Current extraction method using python-pptx"""
        logger.info("Using MO extraction method")
        
        try:
            prs = Presentation(pptx_path)
            
            document = {
                "document_metadata": {
                    "filename": Path(pptx_path).name,
                    "page_count": len(prs.slides),
                    "extracted_at": datetime.now().isoformat(),
                    "extraction_method": "MO"
                },
                "page_de_garde": {},
                "slide_2": {},
                "pages_suivantes": [],
                "page_de_fin": {}
            }
            
            for idx, slide in enumerate(prs.slides, start=1):
                slide_data = {
                    "slide_number": idx,
                    "content": []
                }
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_data["content"].append({
                            "type": "text",
                            "text": shape.text.strip()
                        })
                
                if idx == 1:
                    document["page_de_garde"] = slide_data
                elif idx == 2:
                    document["slide_2"] = slide_data
                elif idx == len(prs.slides):
                    document["page_de_fin"] = slide_data
                else:
                    document["pages_suivantes"].append(slide_data)
            
            if output_path:
                with open(output_path, 'w', encoding='utf-8') as f:
                    json.dump(document, f, indent=2, ensure_ascii=False)
                logger.info(f"Saved extraction to: {output_path}")
            
            return document
            
        except Exception as e:
            logger.error(f"MO extraction failed: {e}")
            raise
    
    def _extract_fida(self, pptx_path: str, output_path: str = None, on_event: Optional[callable] = None) -> Dict[str, Any]:
        """Fida's extraction method using Gemini"""
        logger.info("Using FD (Fida) extraction method with Gemini")
        
        if not self.gemini_api_key:
            logger.warning("GEMINI_API_KEY not set, falling back to MO method")
            return self._extract_mo(pptx_path, output_path)
        
        try:
            from fida import PPTXFinancialExtractor
            
            extractor = PPTXFinancialExtractor(self.gemini_api_key)
            raw_data = extractor.extract_raw_pptx(pptx_path)
            extracted_data = extractor.extract_with_llm_batched(raw_data, on_event=on_event)
            
            if "document_metadata" in extracted_data:
                extracted_data["document_metadata"]["extraction_method"] = "FD"
            
            if output_path:
                with open(output_path, 'w', encoding='utf-8') as f:
                    json.dump(extracted_data, f, indent=2, ensure_ascii=False)
                logger.info(f"Saved FD extraction to: {output_path}")
            
            return extracted_data
            
        except ImportError as e:
            logger.error(f"Import error for fida.py: {e}, falling back to MO method")
            return self._extract_mo(pptx_path, output_path)
        except Exception as e:
            logger.error(f"FD extraction failed: {e}, falling back to MO")
            import traceback
            logger.error(traceback.format_exc())
            return self._extract_mo(pptx_path, output_path)
    
    def _extract_safa(self, pptx_path: str, output_path: str = None) -> Dict[str, Any]:
        """Safa's extraction method using Groq - requires pre-extracted JSON and checklist"""
        logger.info("Using SF (Safa) extraction method with Groq")
        
        if not self.groq_api_key:
            logger.warning("GROQ_API_KEY not set, falling back to MO method")
            return self._extract_mo(pptx_path, output_path)
        
        try:
            # Safa expects pre-extracted JSON, so first do MO extraction
            # Then convert to Safa's expected format
            mo_data = self._extract_mo(pptx_path, None)
            
            # Convert MO format to Safa's expected "elements" format
            elements = []
            
            # Add page_de_garde content
            if mo_data.get("page_de_garde"):
                for item in mo_data["page_de_garde"].get("content", []):
                    elements.append({
                        "page_number": 1,
                        "content": item.get("text", "")
                    })
            
            # Add slide_2 content
            if mo_data.get("slide_2"):
                for item in mo_data["slide_2"].get("content", []):
                    elements.append({
                        "page_number": 2,
                        "content": item.get("text", "")
                    })
            
            # Add pages_suivantes content
            for page in mo_data.get("pages_suivantes", []):
                slide_num = page.get("slide_number", 3)
                for item in page.get("content", []):
                    elements.append({
                        "page_number": slide_num,
                        "content": item.get("text", "")
                    })
            
            # Add page_de_fin content
            if mo_data.get("page_de_fin"):
                page_count = mo_data.get("document_metadata", {}).get("page_count", 6)
                for item in mo_data["page_de_fin"].get("content", []):
                    elements.append({
                        "page_number": page_count,
                        "content": item.get("text", "")
                    })
            
            # Create Safa-compatible input
            safa_input = {
                "filename": mo_data.get("document_metadata", {}).get("filename", ""),
                "elements": elements
            }
            
            # Note: Safa requires a checklist file which we don't have
            # For now, return the converted format with SF marker
            result = mo_data.copy()
            result["document_metadata"]["extraction_method"] = "SF"
            result["safa_elements"] = elements
            
            logger.info(f"SF extraction: Converted {len(elements)} elements")
            
            if output_path:
                with open(output_path, 'w', encoding='utf-8') as f:
                    json.dump(result, f, indent=2, ensure_ascii=False)
                logger.info(f"Saved SF extraction to: {output_path}")
            
            return result
            
        except Exception as e:
            logger.error(f"SF extraction failed: {e}, falling back to MO")
            import traceback
            logger.error(traceback.format_exc())
            return self._extract_mo(pptx_path, output_path)
    
    def _extract_slim(self, pptx_path: str, output_path: str = None, parallel_workers: int = 4) -> Dict[str, Any]:
        """Slim's extraction method using TokenFactory with parallel execution"""
        # Clamp workers to valid range (1-8)
        parallel_workers = max(1, min(8, parallel_workers))
        logger.info(f"Using SL (Slim) extraction method with TokenFactory ({parallel_workers} workers)")
        
        if not self.tokenfactory_api_key:
            logger.warning("TOKENFACTORY_API_KEY not set, falling back to MO method")
            return self._extract_mo(pptx_path, output_path)
        
        try:
            from slim import DocumentParser, extract_text_from_pptx
            
            # Extract text from PPTX
            document_text = extract_text_from_pptx(pptx_path)
            
            if not document_text:
                logger.warning("No text extracted from PPTX, falling back to MO")
                return self._extract_mo(pptx_path, output_path)
            
            # Parse with Slim's parallel extraction using specified workers
            parser = DocumentParser(num_workers=parallel_workers)
            slim_result = parser.parse_document_parallel(document_text)
            
            # Convert Slim output to our standard format
            result = self._convert_slim_to_standard(slim_result, pptx_path)
            
            if output_path:
                with open(output_path, 'w', encoding='utf-8') as f:
                    json.dump(result, f, indent=2, ensure_ascii=False)
                logger.info(f"Saved SL extraction to: {output_path}")
            
            return result
            
        except ImportError as e:
            logger.error(f"Import error for slim.py: {e}, falling back to MO method")
            return self._extract_mo(pptx_path, output_path)
        except Exception as e:
            logger.error(f"SL extraction failed: {e}, falling back to MO")
            import traceback
            logger.error(traceback.format_exc())
            return self._extract_mo(pptx_path, output_path)
    
    def _convert_slim_to_standard(self, slim_result: Dict, pptx_path: str) -> Dict[str, Any]:
        """Convert Slim's output format to our standard format"""
        prs = Presentation(pptx_path)
        page_count = len(prs.slides)
        
        # Build standard format from Slim result
        document = {
            "document_metadata": {
                "filename": Path(pptx_path).name,
                "page_count": page_count,
                "extracted_at": datetime.now().isoformat(),
                "extraction_method": "SL",
                **slim_result.get("document_metadata", {})
            },
            "page_de_garde": {"slide_number": 1, "content": []},
            "slide_2": {"slide_number": 2, "content": []},
            "pages_suivantes": [],
            "page_de_fin": {"slide_number": page_count, "content": []},
            # Keep Slim's rich data
            "claims": slim_result.get("claims", []),
            "disclaimers_and_warnings": slim_result.get("disclaimers_and_warnings", []),
            "key_data_points": slim_result.get("key_data_points", {}),
            "compliance_metadata": slim_result.get("compliance_metadata", {}),
            "performance_metrics": slim_result.get("performance_metrics", {})
        }
        
        # Extract slide content using MO method for structure
        for idx, slide in enumerate(prs.slides, start=1):
            slide_content = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append({
                        "type": "text",
                        "text": shape.text.strip()
                    })
            
            if idx == 1:
                document["page_de_garde"]["content"] = slide_content
            elif idx == 2:
                document["slide_2"]["content"] = slide_content
            elif idx == page_count:
                document["page_de_fin"]["content"] = slide_content
            else:
                document["pages_suivantes"].append({
                    "slide_number": idx,
                    "content": slide_content
                })
        
        return document


# Singleton instance
extraction_manager = ExtractionManager()
