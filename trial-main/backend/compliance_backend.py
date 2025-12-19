#!/usr/bin/env python3
"""
Compliance Backend Wrapper
Bridges the ComplianceOrchestrator with the FastAPI interface
"""

import json
import sys
import shutil
from pathlib import Path
from typing import Dict, Any, Optional
from datetime import datetime
import time

# Import path utilities
from path_utils import (
    DOCUMENTS_DIR, RULES_DIR, 
    get_document_file, get_rule_file
)

# Import the existing orchestrator
from run_all_compliance_checks import ComplianceOrchestrator


class ComplianceBackend:
    """
    Wrapper class that adapts ComplianceOrchestrator for the REST API
    """
    
    def __init__(self, work_dir: str = "./work"):
        """
        Initialize the compliance backend
        
        Args:
            work_dir: Working directory for processing files
        """
        self.work_dir = Path(work_dir)
        self.work_dir.mkdir(parents=True, exist_ok=True)
        
        # Copy database files to work directory if they don't exist locally
        self._ensure_database_files()
        
        self.orchestrator = None
        self.results = None
    
    def _ensure_database_files(self):
        """
        Ensure all required database files are accessible
        Creates symbolic links or copies them to work directory if needed
        """
        # Database files from path-utils
        db_mappings = {
            'disclaimers': (get_document_file('disclaimers'), 'disclaimers.csv'),
            'registration': (get_document_file('registration'), 'registration.csv'),
            'esg_rules': (get_rule_file('esg'), 'esg_rules.json'),
            'general_rules': (get_rule_file('general'), 'general_rules.json'),
            'performance_rules': (get_rule_file('performance'), 'performance_rules.json'),
            'prospectus_rules': (get_rule_file('prospectus'), 'prospectus_rules.json'),
            'structure_rules': (get_rule_file('structure'), 'structure_rules.json'),
            'values_rules': (get_rule_file('values'), 'values_rules.json')
        }
        
        # Try to locate and link database files
        for key, (source_path, dest_name) in db_mappings.items():
            dest_path = Path(dest_name)
            
            # If destination doesn't exist, try to copy from source
            if not dest_path.exists() and source_path.exists():
                print(f"Copying {source_path} to {dest_path}")
                shutil.copy(source_path, dest_path)
    
    def run_full_pipeline(
        self, 
        pptx_path: str, 
        metadata_path: str,
        prospectus_path: Optional[str] = None,
        selected_modules: Optional[list] = None
    ) -> Dict[str, Any]:
        """
        Run the full compliance validation pipeline
        
        Args:
            pptx_path: Path to PowerPoint file
            metadata_path: Path to metadata JSON
            prospectus_path: Optional path to prospectus document
            
        Returns:
            Dictionary with pipeline results
        """
        start_time = time.time()
        
        try:
            # Step 1: Check for existing extraction or extract PowerPoint to JSON
            # Look for pre-extracted JSON in the same directory as the PPTX
            pptx_dir = Path(pptx_path).parent
            existing_extraction = pptx_dir / "extracted_document.json"
            
            if existing_extraction.exists():
                print("📄 Step 1/4: Using existing extraction (already done in background)...")
                document_json_path = existing_extraction
            else:
                print("📄 Step 1/4: Extracting PowerPoint content...")
                document_json_path = self._extract_pptx_to_json(pptx_path)
            
            # Step 2: Validate prospectus path - look in documents folder if not provided
            if prospectus_path is None or not Path(prospectus_path).exists():
                # Try to find prospectus in the documents folder
                backend_dir = Path(__file__).parent
                documents_dir = backend_dir.parent / "documents"
                
                # Look for prospectus.docx in documents folder
                default_prospectus = documents_dir / "prospectus.docx"
                
                if default_prospectus.exists():
                    print(f"📄 Using default prospectus from: {default_prospectus}")
                    prospectus_path = str(default_prospectus)
                else:
                    # Also try to find any .docx file in documents folder
                    docx_files = list(documents_dir.glob("*.docx")) if documents_dir.exists() else []
                    if docx_files:
                        prospectus_path = str(docx_files[0])
                        print(f"📄 Using prospectus found: {prospectus_path}")
                    else:
                        print("⚠️  No prospectus found, Prospectus module may fail")
                        # Create a dummy prospectus as last resort
                        prospectus_path = self.work_dir / "dummy_prospectus.docx"
                        if not Path(prospectus_path).exists():
                            from docx import Document
                            doc = Document()
                            doc.add_paragraph("Placeholder prospectus document")
                            doc.save(str(prospectus_path))
                            print(f"📄 Created placeholder prospectus at: {prospectus_path}")
            
            # Step 3: Run compliance orchestrator
            print("🔍 Step 2/4: Running compliance validation...")
            self.orchestrator = ComplianceOrchestrator(
                document_path=str(document_json_path),
                prospectus_path=str(prospectus_path),
                metadata_path=str(metadata_path)
            )
            
            # Validate prerequisites
            if not self.orchestrator.validate_prerequisites():
                raise Exception("Prerequisites validation failed")
            
            # Run all modules or selected modules
            print(f"🔄 Running modules: {selected_modules if selected_modules else 'ALL'}")
            try:
                if selected_modules:
                    self.orchestrator.run_selected_modules(selected_modules)
                else:
                    self.orchestrator.run_all_modules()
                print("✅ Module execution completed")
            except Exception as module_error:
                print(f"⚠️ Module execution error: {module_error}")
                import traceback
                traceback.print_exc()
            
            # Step 4: Consolidate violations
            print("📊 Step 3/4: Consolidating violations...")
            violations = self.orchestrator.consolidate_violations()
            
            # Step 5: Generate reports
            print("📝 Step 4/4: Generating reports...")
            report_text, json_output = self.orchestrator.generate_master_report()
            
            # Calculate statistics
            total_violations = len(violations)
            critical_violations = sum(1 for v in violations if v.severity == 'critical')
            major_violations = sum(1 for v in violations if v.severity == 'major')
            minor_violations = sum(1 for v in violations if v.severity == 'minor')
            
            duration = time.time() - start_time
            
            # Prepare result
            self.results = {
                'success': True,
                'duration_seconds': duration,
                'summary': {
                    'total_violations': total_violations,
                    'critical_violations': critical_violations,
                    'major_violations': major_violations,
                    'minor_violations': minor_violations,
                    'modules_executed': len(self.orchestrator.MODULES),
                    'pages_with_violations': len(set(v.page_number for v in violations if v.page_number > 0))
                },
                'violations': [
                    {
                        'rule_id': v.rule_id,
                        'module': v.module,
                        'severity': v.severity,
                        'page_number': v.page_number,
                        'location': v.location,
                        'exact_phrase': v.exact_phrase,
                        'violation_comment': v.violation_comment,
                        'required_action': v.required_action,
                        'rule_category': v.rule_category
                    }
                    for v in violations
                ],
                'module_results': {
                    log['module']: {
                        'success': log['success'],
                        'violation_count': log.get('violation_count', 0)
                    }
                    for log in self.orchestrator.execution_log
                }
            }
            
            print(f"✅ Pipeline completed in {duration:.2f} seconds")
            print(f"   Total violations: {total_violations}")
            print(f"   Critical: {critical_violations}")
            print(f"   Major: {major_violations}")
            print(f"   Minor: {minor_violations}")
            
            return self.results
            
        except Exception as e:
            duration = time.time() - start_time
            print(f"❌ Pipeline failed: {str(e)}")
            import traceback
            traceback.print_exc()
            
            return {
                'success': False,
                'duration_seconds': duration,
                'error': str(e),
                'summary': {
                    'total_violations': 0,
                    'critical_violations': 0,
                    'major_violations': 0,
                    'minor_violations': 0
                }
            }
    
    def _extract_pptx_to_json(self, pptx_path: str) -> Path:
        """
        Extract PowerPoint content to JSON format
        
        Args:
            pptx_path: Path to PowerPoint file
            
        Returns:
            Path to generated JSON file
        """
        from pptx import Presentation
        
        output_path = self.work_dir / "extracted_document.json"
        
        try:
            prs = Presentation(pptx_path)
            
            # Build document structure
            document = {
                "document_metadata": {
                    "filename": Path(pptx_path).name,
                    "page_count": len(prs.slides),
                    "extracted_at": datetime.now().isoformat()
                },
                "page_de_garde": {},
                "pages_suivantes": [],
                "page_de_fin": {}
            }
            
            # Extract each slide
            for idx, slide in enumerate(prs.slides, start=1):
                slide_data = {
                    "slide_number": idx,
                    "content": []
                }
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_data["content"].append({
                            "type": "text",
                            "text": shape.text.strip()
                        })
                
                # Classify slide
                if idx == 1:
                    document["page_de_garde"] = slide_data
                elif idx == len(prs.slides):
                    document["page_de_fin"] = slide_data
                else:
                    document["pages_suivantes"].append(slide_data)
            
            # Save to JSON
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(document, f, indent=2, ensure_ascii=False)
            
            print(f"✅ Extracted {len(prs.slides)} slides to {output_path}")
            return output_path
            
        except Exception as e:
            print(f"❌ PowerPoint extraction failed: {e}")
            raise
    
    def save_pipeline_result(self, results: Dict[str, Any], output_path: str):
        """
        Save pipeline results to JSON file
        
        Args:
            results: Pipeline results dictionary
            output_path: Path to save results
        """
        try:
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(results, f, indent=2, ensure_ascii=False)
            print(f"✅ Results saved to {output_path}")
        except Exception as e:
            print(f"❌ Failed to save results: {e}")
    
    def get_summary(self) -> Dict[str, Any]:
        """
        Get summary of last pipeline execution
        
        Returns:
            Summary dictionary
        """
        if self.results is None:
            return {
                'success': False,
                'message': 'No pipeline has been executed yet'
            }
        
        return self.results.get('summary', {})
    
    def get_violations(self) -> list:
        """
        Get list of all violations from last pipeline execution
        
        Returns:
            List of violation dictionaries
        """
        if self.results is None:
            return []
        
        return self.results.get('violations', [])


# CLI interface for testing
def main():
    """CLI interface for testing the backend"""
    if len(sys.argv) < 3:
        print("Usage: python compliance_backend.py <pptx_file> <metadata_file> [prospectus_file]")
        print("\nExample:")
        print("  python compliance_backend.py presentation.pptx metadata.json prospectus.docx")
        sys.exit(1)
    
    pptx_path = sys.argv[1]
    metadata_path = sys.argv[2]
    prospectus_path = sys.argv[3] if len(sys.argv) > 3 else None
    
    print("🚀 Starting Compliance Backend Test")
    print("="*80)
    
    backend = ComplianceBackend(work_dir="./test_work")
    
    results = backend.run_full_pipeline(
        pptx_path=pptx_path,
        metadata_path=metadata_path,
        prospectus_path=prospectus_path
    )
    
    print("\n" + "="*80)
    print("📊 RESULTS SUMMARY")
    print("="*80)
    print(json.dumps(results.get('summary', {}), indent=2))
    
    # Save results
    backend.save_pipeline_result(results, "./test_results.json")
    print("\n✅ Test complete! Check test_results.json for full output")


if __name__ == "__main__":
    main()