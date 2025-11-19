import os
import json
import re
import httpx
from pathlib import Path
from pptx import Presentation
from openai import OpenAI
from typing import Dict, List, Any
from datetime import datetime

class PPTXFinancialExtractor:
    """
    Extracteur de données brutes pour présentations financières
    Version améliorée - extraction exhaustive sans interprétation
    """
    
    def __init__(self, api_key: str):
        http_client = httpx.Client(verify=False)
        self.client = OpenAI(
            api_key=api_key,
            base_url="https://tokenfactory.esprit.tn/api",
            http_client=http_client
        )
        self.model = "hosted_vllm/Llama-3.1-70B-Instruct"
        self.max_slides_per_call = 8  # Réduit pour plus de précision
    
    def extract_raw_pptx(self, pptx_path: str) -> Dict[str, Any]:
        """Extraction brute EXHAUSTIVE du contenu PPTX"""
        prs = Presentation(pptx_path)
        
        slides_content = []
        for idx, slide in enumerate(prs.slides, 1):
            slide_data = {
                "slide_number": idx,
                "texts": [],
                "tables": [],
                "notes": "",
                "has_title": False,
                "shape_count": len(slide.shapes),
                "all_text_raw": ""  # Nouveau: tout le texte brut
            }
            
            all_texts = []
            
            # Extraction EXHAUSTIVE des textes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text_content = shape.text.strip()
                    all_texts.append(text_content)
                    
                    text_info = {
                        "content": text_content,
                        "font_size": None,
                        "font_bold": False,
                        "is_title": False,
                        "position": "body",
                        "shape_type": str(shape.shape_type) if hasattr(shape, 'shape_type') else "unknown"
                    }
                    
                    # Détection du titre
                    if hasattr(shape, "is_placeholder") and shape.is_placeholder:
                        if shape.placeholder_format.type == 1:  # TITLE
                            text_info["is_title"] = True
                            slide_data["has_title"] = True
                    
                    # Formatage détaillé
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.bold:
                                    text_info["font_bold"] = True
                                if run.font.size:
                                    text_info["font_size"] = run.font.size.pt
                                break
                            break
                    
                    # Position
                    if hasattr(shape, "top") and hasattr(shape, "left"):
                        if shape.top < prs.slide_height * 0.15:
                            text_info["position"] = "header"
                        elif shape.top > prs.slide_height * 0.85:
                            text_info["position"] = "footer"
                        
                        text_info["top"] = int(shape.top)
                        text_info["left"] = int(shape.left)
                    
                    slide_data["texts"].append(text_info)
                
                # Tableaux - extraction COMPLÈTE
                if shape.has_table:
                    table_content = {
                        "rows": [],
                        "row_count": len(shape.table.rows),
                        "col_count": len(shape.table.columns)
                    }
                    for row in shape.table.rows:
                        row_data = []
                        for cell in row.cells:
                            cell_text = cell.text.strip()
                            all_texts.append(cell_text)
                            row_data.append(cell_text)
                        table_content["rows"].append(row_data)
                    slide_data["tables"].append(table_content)
            
            # Stocker tout le texte brut
            slide_data["all_text_raw"] = "\n".join(all_texts)
            
            # Notes
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                slide_data["notes"] = notes_text
                all_texts.append(f"[NOTES: {notes_text}]")
            
            slides_content.append(slide_data)
        
        return {
            "filename": Path(pptx_path).name,
            "total_slides": len(prs.slides),
            "slides": slides_content
        }
    
    def extract_with_llm_batched(self, raw_data: Dict) -> Dict[str, Any]:
        """Extraction structurée EXHAUSTIVE sans interprétation"""
        
        total_slides = raw_data['total_slides']
        
        result = {
            "document_metadata": {
                "document_type": "fund_presentation",
                "document_name": "",
                "creation_date": datetime.now().strftime("%Y-%m-%d"),
                "validation_date": "",
                "version": "1.0",
                "client_type": "",
                "language": "",
                "country": "",
                "page_count": total_slides,
                "validated_by": "",
                "fund_esg_classification": "",
                "fund_isin": "",
                "fund_type": ""
            },
            "page_de_garde": {},
            "slide_2": {},
            "pages_suivantes": [],
            "page_de_fin": {}
        }
        
        # 1. Métadonnées du document
        print("📋 Extraction des métadonnées...")
        metadata = self._extract_document_metadata(raw_data['slides'][:3])
        result["document_metadata"].update(metadata)
        
        # 2. Page de garde
        print("📄 Extraction Page de Garde (Slide 1)...")
        result["page_de_garde"] = self._extract_page_garde(raw_data['slides'][0])
        
        # 3. Slide 2
        if total_slides >= 2:
            print("📄 Extraction Slide 2...")
            result["slide_2"] = self._extract_slide_2(raw_data['slides'][1])
        
        # 4. Pages suivantes - TRAITEMENT SLIDE PAR SLIDE
        if total_slides > 3:
            print(f"📄 Extraction Pages Suivantes (Slides 3-{total_slides-1})...")
            middle_slides = raw_data['slides'][2:-1]
            
            for i, slide in enumerate(middle_slides):
                slide_num = i + 3
                print(f"   Processing slide {slide_num}/{total_slides-1}...", end=" ")
                
                slide_result = self._extract_single_slide(slide, slide_num)
                
                if slide_result and slide_result.get("content"):
                    result["pages_suivantes"].append(slide_result)
                    print("✓")
                else:
                    print("⚠️ (minimal content)")
        
        # 5. Page de fin
        if total_slides > 2:
            print(f"📄 Extraction Page de Fin (Slide {total_slides})...")
            result["page_de_fin"] = self._extract_page_fin(raw_data['slides'][-1], total_slides)
        
        return result
    
    def _extract_single_slide(self, slide: Dict, slide_num: int) -> Dict:
        """Extraction EXHAUSTIVE d'une slide individuelle"""
        
        # Vérification rapide : slide vide ou titre uniquement
        all_text = slide.get("all_text_raw", "").strip()
        if len(all_text) < 20 and not slide.get("tables"):
            # Slide de transition simple
            return {
                "slide_number": slide_num,
                "slide_title": all_text if all_text else f"Slide {slide_num}",
                "content": {
                    "main_text": all_text,
                    "slide_type": "transition"
                }
            }
        
        # Préparer les données complètes
        slide_json = json.dumps({
            "slide_number": slide_num,
            "all_text_raw": slide.get("all_text_raw", ""),
            "texts": slide.get("texts", []),
            "tables": slide.get("tables", []),
            "notes": slide.get("notes", ""),
            "has_title": slide.get("has_title", False),
            "shape_count": slide.get("shape_count", 0)
        }, ensure_ascii=False, indent=2)
        
        prompt = f"""Extract ALL data from this slide. Return EVERYTHING you find.

COMPLETE SLIDE DATA:
{slide_json[:8000]}

Return JSON in this EXACT format:
{{
  "slide_number": {slide_num},
  "slide_title": "Extract title if present, otherwise first major text",
  "content": {{
    // IMPORTANT: Only include fields that have actual data. Omit empty fields entirely.
    // Extract EVERYTHING found on this slide into appropriate fields:
    
    "main_text": "Main body text if present",
    "bullet_points": ["List all bullet points exactly as written"],
    
    // Team slides:
    "team_members": [{{"name": "", "role": "", "experience": ""}}],
    
    // Objective/Strategy slides:
    "objective": {{"text": "", "font_size": null, "font_bold": false}},
    "strategy": {{"text": "", "font_size": null, "font_bold": false}},
    
    // Holdings slides:
    "holdings_intro": {{"text": ""}},
    "top_holdings": [{{"name": "", "weight": ""}}],
    
    // Performance slides:
    "performance_data": {{"periods": [], "values": [], "benchmark": "", "disclaimer": ""}},
    
    // SRI/Risk slides:
    "sri_indicator": {{"value": "", "text": "", "font_size": null, "font_bold": false}},
    "risk_description": "",
    
    // ESG slides:
    "esg_approach": {{"text": "", "approach_type": ""}},
    
    // Glossary slides:
    "glossary_terms": [{{
      "term": "",
      "definition": "",
      "font_size": null,
      "font_bold": false
    }}],
    
    // Disclaimers (ALWAYS extract if present):
    "disclaimers": [{{"text": "", "font_size": null, "font_bold": false, "position": ""}}],
    
    // Metadata (extract if present):
    "tables": [/* Copy ALL table data exactly as shown */],
    "charts_description": "Describe any charts/graphs present",
    "sources": ["List ALL sources mentioned"],
    "footnotes": ["List ALL footnotes"],
    "company_mentions": ["List ALL company names mentioned"],
    "numbers_data": ["Extract ALL numbers with their context"],
    "dates": ["Extract ALL dates mentioned"],
    "percentages": ["Extract ALL percentages with context"],
    
    "additional_text": "ANY other text not captured above"
  }}
}}

CRITICAL RULES:
1. Extract EVERY piece of text visible
2. Copy text VERBATIM - no paraphrasing
3. If tables exist, copy them COMPLETELY
4. Extract numbers WITH context (e.g., "AUM: 5M€")
5. Omit fields with no data - don't include empty strings or empty arrays
6. DO NOT add markdown code fences like ```json
7. Return ONLY the JSON object

Return ONLY valid JSON without any markdown formatting."""

        return self._call_llm(prompt)
    
    def _extract_document_metadata(self, first_slides: List[Dict]) -> Dict:
        """Extraction métadonnées - version inchangée"""
        
        all_text = ""
        for slide in first_slides:
            all_text += slide.get("all_text_raw", "") + "\n"
            for text_obj in slide.get("texts", []):
                all_text += text_obj["content"] + "\n"
        
        prompt = f"""Extract document metadata. Return ONLY factual data found.

TEXT CONTENT:
{all_text[:3500]}

Return JSON:
{{
  "document_name": "Exact fund name",
  "validation_date": "Date in YYYY-MM-DD if present",
  "client_type": "retail/professional/institutional if mentioned",
  "language": "Language code (en, fr, de, etc.)",
  "country": "Country code if mentioned",
  "validated_by": "Validator name if mentioned",
  "fund_esg_classification": "ESG classification if mentioned",
  "fund_isin": "ISIN code if present",
  "fund_type": "UCITS/AIF/ETF/other"
}}

RULES:
- Extract ONLY explicit information
- Empty string "" if not found
- Exact wording from document

Return ONLY valid JSON."""

        return self._call_llm(prompt)
    
    def _extract_page_garde(self, slide: Dict) -> Dict:
        """Extraction page de garde - version inchangée mais améliorée"""
        
        slide_json = json.dumps({
            "all_text_raw": slide.get("all_text_raw", ""),
            "texts": slide.get("texts", []),
            "tables": slide.get("tables", [])
        }, ensure_ascii=False, indent=2)
        
        prompt = f"""Extract cover page information. Extract EVERYTHING present.

COMPLETE SLIDE DATA:
{slide_json[:5000]}

Return JSON:
{{
  "slide_number": 1,
  "slide_title": "Title if present",
  "content": {{
    "fund_name": "Exact fund name",
    "management_company": "Exact company name",
    "promotional_document_mention": "Exact promotional mention",
    "target_audience": "Target audience if stated",
    "date": "Date as written",
    "tagline": "Tagline/slogan if present",
    "disclaimer_text": "Disclaimer if on this slide",
    "legal_notice_sgp": "Legal notice if present",
    "additional_text": "Any other text present"
  }}
}}

RULES:
- Extract ALL text present
- VERBATIM extraction only
- Empty string "" for missing data

Return ONLY valid JSON."""

        return self._call_llm(prompt)
    
    def _extract_slide_2(self, slide: Dict) -> Dict:
        """Extraction slide 2 - version améliorée"""
        
        slide_json = json.dumps({
            "all_text_raw": slide.get("all_text_raw", ""),
            "texts": slide.get("texts", []),
            "tables": slide.get("tables", [])
        }, ensure_ascii=False, indent=2)
        
        prompt = f"""Extract slide 2 information EXHAUSTIVELY.

COMPLETE SLIDE DATA:
{slide_json[:6000]}

Return JSON:
{{
  "slide_number": 2,
  "slide_title": "Title",
  "content": {{
    "standard_disclaimer_retail": {{
      "text": "COMPLETE disclaimer text",
      "font_size": null,
      "font_bold": false,
      "visible": true
    }},
    "risk_profile": {{
      "text": "Risk text",
      "font_size": null,
      "font_bold": false,
      "visible": true
    }},
    "capital_at_risk": {{
      "text": "Capital at risk mention",
      "font_size": null,
      "font_bold": false,
      "visible": true,
      "position": "footer/header/body"
    }},
    "distribution_countries": ["ALL countries mentioned"],
    "performance_disclaimer": "Performance disclaimer",
    "team_change_mention": "Team change if mentioned",
    "additional_disclaimers": ["ANY other disclaimers"],
    "all_risks_listed": ["Extract ALL risks mentioned"],
    "tables": [/* Include if present */]
  }}
}}

RULES:
- Extract EVERY disclaimer
- List ALL countries
- List ALL risks
- VERBATIM text only

Return ONLY valid JSON."""

        return self._call_llm(prompt)
    
    def _extract_page_fin(self, slide: Dict, slide_num: int) -> Dict:
        """Extraction page de fin - version améliorée"""
        
        slide_json = json.dumps({
            "slide_number": slide_num,
            "all_text_raw": slide.get("all_text_raw", ""),
            "texts": slide.get("texts", []),
            "tables": slide.get("tables", [])
        }, ensure_ascii=False, indent=2)
        
        prompt = f"""Extract final page COMPLETELY.

COMPLETE SLIDE DATA:
{slide_json[:6000]}

Return JSON:
{{
  "slide_number": {slide_num},
  "slide_title": "Title",
  "content": {{
    "legal_mention_sgp": {{
      "text": "Complete legal text",
      "font_size": null,
      "font_bold": false,
      "visible": true
    }},
    "contact_info": {{
      "management_company": "",
      "address": "Complete address",
      "phone": "",
      "email": "",
      "website": ""
    }},
    "fund_characteristics": {{
      "structure": "",
      "fees": "",
      "minimum_investment": "",
      "benchmark": "",
      "srri": "",
      "aum": ""
    }},
    "validator": {{
      "name": "",
      "role": "",
      "date": ""
    }},
    "document_date": "",
    "additional_disclaimers": ["ALL disclaimers"],
    "tables": [/* ALL tables */],
    "additional_text": "Any other text"
  }}
}}

RULES:
- Extract EVERYTHING
- Complete addresses
- ALL contact info
- VERBATIM text

Return ONLY valid JSON."""

        return self._call_llm(prompt)
    
    def _call_llm(self, prompt: str) -> Dict:
        """Appel API avec gestion d'erreurs robuste"""
        max_retries = 3
        
        for attempt in range(max_retries):
            try:
                response = self.client.chat.completions.create(
                    model=self.model,
                    messages=[
                        {
                            "role": "system",
                            "content": """You are a PRECISE data extraction system.

ABSOLUTE RULES:
1. Extract EVERY piece of information present
2. NEVER interpret, analyze, or add information
3. Copy text EXACTLY as written (verbatim)
4. Return ONLY valid JSON - NO markdown code fences
5. Use empty string "" for missing text
6. Use null for missing numbers
7. OMIT fields that have no data - don't include empty objects or arrays
8. If unsure, extract it anyway - don't omit

CRITICAL: Do NOT wrap your response in ```json or ``` - return raw JSON only.

Remember: MORE is better than LESS. Extract EVERYTHING."""
                        },
                        {
                            "role": "user",
                            "content": prompt
                        }
                    ],
                    temperature=0.0,
                    max_tokens=8000
                )
                
                result_text = response.choices[0].message.content.strip()
                
                # Nettoyage TRÈS agressif
                result_text = re.sub(r'^```json\s*\n?', '', result_text, flags=re.IGNORECASE | re.MULTILINE)
                result_text = re.sub(r'^```\s*\n?', '', result_text, flags=re.MULTILINE)
                result_text = re.sub(r'\n?```\s*', '', result_text, flags=re.MULTILINE)
                result_text = result_text.strip()
                
                # Si commence encore par texte avant JSON, extraire juste le JSON
                if not result_text.startswith('{'):
                    json_start = result_text.find('{')
                    if json_start != -1:
                        result_text = result_text[json_start:]
                
                # Si termine avec du texte après JSON, tronquer
                last_brace = result_text.rfind('}')
                if last_brace != -1 and last_brace < len(result_text) - 1:
                    result_text = result_text[:last_brace+1]
                
                # Parse JSON
                parsed = json.loads(result_text)
                return parsed
                
            except json.JSONDecodeError as e:
                if attempt < max_retries - 1:
                    print(f"⚠️ JSON error (attempt {attempt + 1}/{max_retries})", end="")
                    continue
                else:
                    print(f"❌ JSON error at pos {e.pos}")
                    # Essayer de sauver ce qui est possible
                    try:
                        # Extraire le contenu brut pour analyse
                        return {
                            "slide_number": "unknown",
                            "content": {
                                "extraction_error": f"JSON decode failed at position {e.pos}",
                                "raw_text_sample": result_text[max(0, e.pos-200):min(len(result_text), e.pos+200)]
                            }
                        }
                    except:
                        return {"error": f"JSON decode failed: {str(e)}", "slide_number": "unknown", "content": {}}
            
            except Exception as e:
                if attempt < max_retries - 1:
                    print(f"⚠️ API error: {str(e)[:50]}", end="")
                    continue
                return {"error": str(e), "slide_number": "unknown", "content": {}}
        
        return {"error": "Max retries exceeded", "slide_number": "unknown", "content": {}}
    
    def process(self, pptx_path: str, output_json: str):
        """Pipeline complet d'extraction EXHAUSTIVE"""
        
        print("="*80)
        print("🚀 EXTRACTION PPTX - MODE EXHAUSTIF")
        print("="*80)
        print("\n⚠️  Extraction COMPLÈTE de toutes les données brutes")
        print("   Aucune interprétation - Extraction verbatim uniquement")
        print("="*80)
        
        print("\n📄 Extraction du contenu PPTX...")
        raw_data = self.extract_raw_pptx(pptx_path)
        print(f"✅ {raw_data['total_slides']} slides extraites")
        
        print(f"\n🤖 Extraction structurée avec LLM (slide par slide)...")
        extracted_data = self.extract_with_llm_batched(raw_data)
        
        # Sauvegarder
        with open(output_json, 'w', encoding='utf-8') as f:
            json.dump(extracted_data, f, indent=2, ensure_ascii=False)
        
        print(f"\n✅ Données sauvegardées: {output_json}")
        self._print_summary(extracted_data)
    
    def _print_summary(self, data: Dict):
        """Résumé détaillé de l'extraction"""
        
        print("\n" + "="*80)
        print("📊 RÉSUMÉ DE L'EXTRACTION EXHAUSTIVE")
        print("="*80)
        
        meta = data.get('document_metadata', {})
        print(f"\n📄 Document: {meta.get('document_name', 'N/A')}")
        print(f"📋 Total slides: {meta.get('page_count', 0)}")
        print(f"🌍 Langue/Pays: {meta.get('language', 'N/A')}/{meta.get('country', 'N/A')}")
        print(f"🏷️  Type: {meta.get('fund_type', 'N/A')}")
        print(f"🔢 ISIN: {meta.get('fund_isin', 'N/A')}")
        
        pg = data.get('page_de_garde', {}).get('content', {})
        if pg:
            print(f"\n🟢 Page de garde")
            print(f"   Fonds: {pg.get('fund_name', 'N/A')[:60]}")
            print(f"   Date: {pg.get('date', 'N/A')}")
            print(f"   SGP: {pg.get('management_company', 'N/A')[:50]}")
        
        s2 = data.get('slide_2', {}).get('content', {})
        if s2:
            print(f"\n🟡 Slide 2")
            countries = s2.get('distribution_countries', [])
            print(f"   Pays: {', '.join(countries[:5])}{' ...' if len(countries) > 5 else ''}")
            risks = s2.get('all_risks_listed', [])
            print(f"   Risques listés: {len(risks)}")
        
        pages_suiv = data.get('pages_suivantes', [])
        if pages_suiv:
            print(f"\n🔵 Pages suivantes: {len(pages_suiv)} slides")
            
            # Statistiques détaillées
            slides_with_content = sum(1 for p in pages_suiv if p.get('content') and len(p.get('content', {})) > 2)
            slides_with_tables = sum(1 for p in pages_suiv 
                                    if p.get('content', {}).get('tables'))
            slides_with_team = sum(1 for p in pages_suiv 
                                  if p.get('content', {}).get('team_members'))
            slides_with_glossary = sum(1 for p in pages_suiv 
                                      if p.get('content', {}).get('glossary_terms'))
            slides_with_holdings = sum(1 for p in pages_suiv 
                                      if p.get('content', {}).get('top_holdings'))
            transition_slides = sum(1 for p in pages_suiv 
                                   if p.get('content', {}).get('slide_type') == 'transition')
            
            print(f"   - Avec contenu substantiel: {slides_with_content}/{len(pages_suiv)}")
            print(f"   - Slides de transition: {transition_slides}")
            print(f"   - Avec tableaux: {slides_with_tables}")
            print(f"   - Équipe: {slides_with_team}")
            print(f"   - Holdings: {slides_with_holdings}")
            print(f"   - Glossaire: {slides_with_glossary}")
            
            # Vérifier les slides avec erreurs
            error_slides = [p.get('slide_number') for p in pages_suiv 
                          if 'error' in p or p.get('content', {}).get('extraction_error')]
            if error_slides:
                print(f"   ⚠️ Slides avec erreurs d'extraction: {error_slides}")
            
            # Vérifier les slides minimales (hors transitions)
            minimal_slides = [p.get('slide_number') for p in pages_suiv 
                            if p.get('content') and len(p.get('content', {})) < 2 
                            and p.get('content', {}).get('slide_type') != 'transition']
            if minimal_slides:
                print(f"   ⚠️ Slides avec contenu minimal: {minimal_slides}")
        
        pf = data.get('page_de_fin', {})
        if pf:
            print(f"\n🟣 Page de fin: Slide {pf.get('slide_number', 'N/A')}")
            if 'error' in pf:
                print(f"   ⚠️ Erreur d'extraction: {pf.get('error', 'Unknown')[:60]}")
            else:
                contact = pf.get('content', {}).get('contact_info', {})
                if contact and contact.get('management_company'):
                    print(f"   SGP: {contact.get('management_company')}")
        
        print("\n" + "="*80)
        print("✅ Extraction exhaustive terminée")
        
        # Compter le total d'éléments extraits
        total_tables = len(s2.get('tables', [])) + sum(len(p.get('content', {}).get('tables', [])) for p in pages_suiv)
        total_disclaimers = sum(len(p.get('content', {}).get('disclaimers', [])) for p in pages_suiv)
        
        print(f"\n📈 Statistiques globales:")
        print(f"   - Total tableaux: {total_tables}")
        print(f"   - Total disclaimers: {total_disclaimers}")
        print("="*80)


# ===== UTILISATION =====
if __name__ == "__main__":
    
    # Configuration
    TOKENFACTORY_API_KEY = "sk-f51cc6f011fd41efb18c1db77c7b273f"
    PPTX_FILE = "XXX-PRS-GB-ODDO BHF US Equity Active ETF-20250630_6PN.pptx"
    OUTPUT_JSON = "extracted_data_exhaustive12.json"
    
    # Extraction
    extractor = PPTXFinancialExtractor(api_key=TOKENFACTORY_API_KEY)
    extractor.process(PPTX_FILE, OUTPUT_JSON)
    
    print("\n🎯 Prochaine étape: Analyse de conformité")